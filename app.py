import streamlit as st
import pandas as pd
import io
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches

# Set a professional style for matplotlib plots
plt.style.use('seaborn-v0_8-darkgrid') # Using a darkgrid style for a clean look
sns.set_palette("viridis") # A colorblind-friendly and aesthetically pleasing palette

# Function to analyze CSV content and generate summary
def analyze_csv_data(csv_content: str):
    """
    Analyzes CSV content and generates a summary.

    Args:
        csv_content: A string containing the CSV data.

    Returns:
        A dictionary containing the dataset summary, column information, and the DataFrame.
    """
    try:
        # Read CSV content into a pandas DataFrame
        df = pd.read_csv(io.StringIO(csv_content))
    except Exception as e:
        return {"error": f"Failed to parse CSV: {e}"}

    num_rows, num_cols = df.shape
    missing_values_count = df.isnull().sum().sum()
    total_cells = num_rows * num_cols
    missing_percentage = (missing_values_count / total_cells * 100) if total_cells > 0 else 0

    column_info = {}

    # Analyze each column
    for col_name in df.columns:
        col_data = df[col_name]
        
        # Attempt to convert to numeric, coercing errors to NaN
        numeric_col_data = pd.to_numeric(col_data, errors='coerce').dropna()
        
        # Count missing values for the current column
        col_missing_count = col_data.isnull().sum()

        # Determine if the column is primarily numeric
        # A column is considered numeric if more than 80% of its non-missing values are numeric
        if not numeric_col_data.empty and len(numeric_col_data) / (len(col_data) - col_missing_count) > 0.8:
            column_info[col_name] = {
                'type': 'numeric',
                'min': numeric_col_data.min(),
                'max': numeric_col_data.max(),
                'mean': numeric_col_data.mean(),
                'median': numeric_col_data.median(), # Added median for better central tendency
                'std': numeric_col_data.std(),
                'count': len(numeric_col_data), # Count of non-null numeric values
                'missing_count': col_missing_count
            }
        else:
            # Treat as categorical
            counts = col_data.value_counts().to_dict()
            column_info[col_name] = {
                'type': 'categorical',
                'unique_values': len(counts),
                'counts': counts,
                'missing_count': col_missing_count
            }

    # Generate summary text
    summary = f"The dataset contains {num_rows} rows and {num_cols} columns.\n"
    summary += f"Total missing values across the dataset: {missing_values_count} ({missing_percentage:.2f}% of all cells).\n\n"
    summary += "Column Details:\n"
    for col_name, info in column_info.items():
        summary += f"- **{col_name}**: "
        if info['type'] == 'numeric':
            summary += f"Numeric (Min: {info['min']:.2f}, Max: {info['max']:.2f}, Mean: {info['mean']:.2f}, Median: {info['median']:.2f}, Std Dev: {info['std']:.2f}, Non-null: {info['count']}). Missing: {info['missing_count']} values.\n"
        else:
            summary += f"Categorical ({info['unique_values']} unique values). Missing: {info['missing_count']} values.\n"
    
    return {
        "summary": summary,
        "column_info": column_info,
        "dataframe": df # Return the DataFrame for plotting
    }

# --- Streamlit App ---
st.set_page_config(layout="wide", page_title="Dataset Analyzer")

st.title("📊 Dataset Analyzer")
st.markdown("Upload your CSV file to get an automated overview and visualizations.")

uploaded_file = st.file_uploader("Choose a CSV file", type="csv")

if uploaded_file is not None:
    # Read the uploaded file
    csv_content = uploaded_file.getvalue().decode("utf-8")

    with st.spinner("Analyzing dataset and generating visualizations..."):
        analysis_results = analyze_csv_data(csv_content)

    if "error" in analysis_results:
        st.error(f"Error: {analysis_results['error']}")
    else:
        st.subheader("Dataset Overview")
        st.text(analysis_results["summary"]) # st.text preserves formatting

        st.subheader("Visualizations")
        df = analysis_results["dataframe"]
        column_info = analysis_results["column_info"]

        # List to store chart images for PPT
        chart_images_for_ppt = []

        # Use Streamlit columns for better chart layout
        chart_cols = st.columns(2) # Create two columns for charts

        num_charts_displayed = 0
        chart_idx = 0

        for col_name, info in column_info.items():
            # Use the current column for placing the chart
            with chart_cols[chart_idx % 2]:
                if info['type'] == 'categorical' and info['unique_values'] < 50:
                    chart_title = f'Frequency Distribution of {col_name}'
                    st.write(f"#### {chart_title}")
                    # Order categorical bars by frequency for better readability
                    ordered_counts = pd.Series(info['counts']).sort_values(ascending=False)
                    
                    fig, ax = plt.subplots(figsize=(8, 5)) # Set a consistent figure size
                    sns.barplot(x=ordered_counts.index, y=ordered_counts.values, ax=ax, palette="viridis")
                    ax.set_title(chart_title, fontsize=14)
                    ax.set_xlabel(col_name, fontsize=12)
                    ax.set_ylabel('Count', fontsize=12)
                    plt.xticks(rotation=45, ha='right') # Rotate labels for readability
                    plt.tight_layout() # Adjust layout to prevent labels from overlapping
                    st.pyplot(fig)
                    
                    # Save chart to BytesIO for PPT
                    buf = io.BytesIO()
                    fig.savefig(buf, format="png", bbox_inches='tight')
                    buf.seek(0)
                    chart_images_for_ppt.append({"title": chart_title, "image_buffer": buf})
                    plt.close(fig) # Close the figure to free memory
                    num_charts_displayed += 1
                    chart_idx += 1
                elif info['type'] == 'numeric':
                    if info['count'] > 10: # Only plot if sufficient data points
                        chart_title = f'Distribution of {col_name}'
                        st.write(f"#### {chart_title}")
                        fig, ax = plt.subplots(figsize=(8, 5)) # Set a consistent figure size
                        sns.histplot(df[col_name].dropna(), kde=True, ax=ax, palette="viridis")
                        ax.set_title(chart_title, fontsize=14)
                        ax.set_xlabel(col_name, fontsize=12)
                        ax.set_ylabel('Density / Count', fontsize=12)
                        plt.tight_layout()
                        st.pyplot(fig)

                        # Save chart to BytesIO for PPT
                        buf = io.BytesIO()
                        fig.savefig(buf, format="png", bbox_inches='tight')
                        buf.seek(0)
                        chart_images_for_ppt.append({"title": chart_title, "image_buffer": buf})
                        plt.close(fig) # Close the figure to free memory
                        num_charts_displayed += 1
                        chart_idx += 1
                    else:
                        st.info(f"Not enough data points to plot distribution for '{col_name}' in this column.")

        if num_charts_displayed == 0:
            st.info("No suitable columns found for automatic chart generation (e.g., all columns are text, too many unique values, or insufficient numeric data).")

        st.subheader("Download Section")
        st.markdown("You can download the analyzed dataset or a summary of the analysis for your records.")
        
        # Download original DataFrame as CSV
        csv_download = df.to_csv(index=False).encode('utf-8')
        st.download_button(
            label="Download Analyzed Data (CSV)",
            data=csv_download,
            file_name="analyzed_data.csv",
            mime="text/csv",
            help="Download the full dataset you uploaded, as processed by the analyzer."
        )

        # Download summary text
        summary_download_text = analysis_results["summary"]
        st.download_button(
            label="Download Dataset Overview (Text)",
            data=summary_download_text,
            file_name="dataset_overview.txt",
            mime="text/plain",
            help="Download the textual summary of the dataset's characteristics."
        )

        # --- PowerPoint Generation and Download ---
        st.subheader("Generate Presentation (PPT)")
        if chart_images_for_ppt:
            if st.button("Generate and Download PPT"):
                prs = Presentation()
                # Use a blank slide layout (layout index 6 is typically blank)
                blank_slide_layout = prs.slide_layouts[6] 

                # Add a title slide
                title_slide_layout = prs.slide_layouts[0] # Title slide layout
                slide = prs.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                title.text = "Dataset Analysis Report"
                subtitle.text = "Generated by Streamlit Analyzer"

                # Add slides for each chart
                for chart_data in chart_images_for_ppt:
                    slide = prs.slides.add_slide(blank_slide_layout)
                    # Add title to the slide
                    left = top = width = height = Inches(0.5) # Placeholder values
                    title_shape = slide.shapes.add_textbox(left, top, prs.slide_width - Inches(1), Inches(0.75))
                    tf = title_shape.text_frame
                    tf.text = chart_data["title"]
                    tf.paragraphs[0].font.size = Inches(0.25) # Adjust font size if needed

                    # Calculate position to center the image
                    img_width = Inches(8)
                    img_height = Inches(5)
                    left = (prs.slide_width - img_width) / 2
                    top = (prs.slide_height - img_height) / 2 + Inches(0.5) # Adjust for title

                    slide.shapes.add_picture(chart_data["image_buffer"], left, top, img_width, img_height)

                # Save presentation to a BytesIO object
                ppt_buffer = io.BytesIO()
                prs.save(ppt_buffer)
                ppt_buffer.seek(0)

                st.download_button(
                    label="Click to Download PPT",
                    data=ppt_buffer,
                    file_name="dataset_analysis_report.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    help="Download a PowerPoint presentation with all generated charts."
                )
                st.success("PPT generated successfully!")
        else:
            st.info("Upload a dataset to generate charts and the PPT.")


        st.subheader("Context: Leveraging Data for Organizational Development")
        st.markdown("""
        After analyzing the charts and dataset overview, here's how you can translate these insights into actionable strategies for your organization:

        **1. Identify Trends and Patterns:**
        * **Categorical Distributions:** Look at bar charts for categorical data (e.g., 'City', 'Occupation'). Are there dominant categories? Are there unexpected distributions? This can inform resource allocation, targeted marketing, or identifying underserved segments. For instance, a high concentration in one city might suggest a need for more localized support or expansion opportunities elsewhere.
        * **Numeric Distributions:** Histograms and summary statistics for numeric data (e.g., 'Age', 'Salary', 'Experience') reveal the spread and central tendencies. Are salaries clustered around a certain range? Is age distribution skewed? This can help in workforce planning, compensation strategies, or identifying training needs.

        **2. Spot Anomalies and Outliers:**
        * Unusual spikes or dips in charts, or extreme min/max values in numeric summaries, could indicate data entry errors, rare events, or critical issues. Investigating these can uncover fraud, system glitches, or unique business opportunities.

        **3. Assess Data Quality:**
        * The "Missing values" count is crucial. High percentages of missing data in key columns can impact the reliability of your analysis. Consider strategies for data imputation or collecting more complete data.

        **4. Formulate Hypotheses:**
        * Based on visual patterns, start asking "why" questions. Why is one category so much larger than others? Why is there a bimodal distribution in a numeric column? These hypotheses can guide further, more targeted analysis.

        **5. Drive Strategic Decisions:**
        * **Resource Allocation:** If certain product categories are consistently popular (from categorical charts), allocate more resources to their development or marketing.
        * **Operational Efficiency:** Identify bottlenecks or areas of inefficiency by analyzing process-related data.
        * **Customer Understanding:** Understand customer demographics, preferences, or behavior patterns from relevant columns to tailor products or services.
        * **Risk Management:** Detect unusual patterns that might indicate potential risks (e.g., fraud, equipment failure).

        **6. Plan Further Deep Dives:**
        * This initial analysis is exploratory. Consider more advanced techniques:
            * **Correlation Analysis:** Explore relationships between numeric variables (e.g., using scatter plots or correlation matrices).
            * **Segmentation:** Group your data based on certain characteristics to understand different customer segments or employee groups.
            * **Predictive Modeling:** If you have a target variable (e.g., 'Churn', 'Sales'), you might consider building machine learning models to predict future outcomes.
            * **Time Series Analysis:** If your data has a time component, analyze trends over time.

        By systematically interpreting these initial insights, organizations can move from raw data to informed decisions, fostering continuous improvement and strategic growth.
        """)
        
        st.subheader("Creating a Professional Presentation (PPT) - Manual Steps")
        st.markdown("""
        While you can now download an automated PPT, for more customized and polished presentations, you might still prefer to manually create one using the following steps:

        1.  **Download Charts as Images:** For each chart displayed above, right-click on the chart image in your browser and select "Save Image As..." to save it as a PNG file. These high-quality images can be directly inserted into your PowerPoint slides.
        2.  **Copy Dataset Overview:** Copy the "Dataset Overview" text directly from the Streamlit app or download it using the "Download Dataset Overview (Text)" button. This text provides a concise summary of your data's structure and characteristics, perfect for an introductory slide.
        3.  **Utilize "Context for Organizational Development":** The "Context: Leveraging Data for Organizational Development" section provides structured bullet points and explanations. You can copy and paste these points directly into your slides, expanding on them with specific examples from your dataset.
        4.  **Structure Your Presentation:**
            * **Title Slide:** Your presentation title and company logo.
            * **Introduction:** Briefly describe the dataset and its purpose.
            * **Dataset Overview:** Use the downloaded text summary.
            * **Key Visualizations:** Dedicate slides to each important chart, explaining what it shows and its immediate implications.
            * **Insights & Recommendations:** Use the "Context" section to discuss trends, anomalies, and actionable strategies. Tailor these to your specific organizational goals.
            * **Next Steps/Further Analysis:** Outline what deeper dives or additional data collection might be beneficial.
            * **Conclusion:** Summarize key findings and call to action.
        5.  **Design and Branding:** Apply your organization's branding, colors, and fonts to the PowerPoint template to maintain a professional and consistent look.
        6.  **Tell a Story:** The most effective presentations tell a compelling story with data. Connect your charts and insights to a clear narrative that addresses a business question or problem.
        """)

else:
    st.info("Please upload a CSV file to begin analysis.")

st.markdown("---")
st.markdown("Developed with Streamlit")
